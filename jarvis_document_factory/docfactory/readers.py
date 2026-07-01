"""Format readers: extract text and embedded-image info from rendered files.

Shared by the orchestrator (cross-format consistency) and the Output_Gate.
PDF and DOCX are read with pypdf and python-docx; PPTX with python-pptx.
"""
from __future__ import annotations


def read_pdf_pages(path: str) -> list[str]:
    from pypdf import PdfReader

    reader = PdfReader(path)
    return [(p.extract_text() or "") for p in reader.pages]


def read_docx_paragraph_texts(path: str) -> list[str]:
    import docx

    d = docx.Document(path)
    texts = [p.text for p in d.paragraphs]
    for tbl in d.tables:
        for row in tbl.rows:
            for cell in row.cells:
                texts.append(cell.text)
    for section in d.sections:
        for p in section.footer.paragraphs:
            texts.append(p.text)
    return texts


def read_pptx_slide_texts(path: str) -> list[str]:
    from pptx import Presentation

    prs = Presentation(path)
    slides = []
    for slide in prs.slides:
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                parts.append(shape.text_frame.text)
        slides.append("\n".join(parts))
    return slides


def read_text(fmt: str, path: str) -> str:
    if fmt == "pdf":
        return "\n".join(read_pdf_pages(path))
    if fmt == "docx":
        return "\n".join(read_docx_paragraph_texts(path))
    if fmt == "pptx":
        return "\n".join(read_pptx_slide_texts(path))
    raise ValueError("unknown format: %r" % fmt)


# ------------------------------------------------------------------ images
def count_pdf_images(path: str) -> int:
    """Count DISTINCT embedded image XObjects.

    WeasyPrint shares one image XObject across every page's resource dict, so a
    single embedded image can appear in multiple pages. Dedupe by the indirect
    object reference so the count reflects distinct images, not page references.
    """
    from pypdf import PdfReader

    reader = PdfReader(path)
    seen = set()
    for page in reader.pages:
        try:
            res = page.get("/Resources")
            xo = res.get("/XObject") if res else None
            if xo:
                xo = xo.get_object()
                for value in xo.values():
                    obj = value.get_object()
                    if obj.get("/Subtype") == "/Image":
                        ref = getattr(value, "indirect_reference", None)
                        key = (ref.idnum, ref.generation) if ref is not None else id(obj)
                        seen.add(key)
        except Exception:
            pass
    return len(seen)


def count_docx_images(path: str) -> int:
    import docx

    return len(docx.Document(path).inline_shapes)


def count_pptx_images(path: str) -> int:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(path)
    n = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                n += 1
    return n


def count_images(fmt: str, path: str) -> int:
    if fmt == "pdf":
        return count_pdf_images(path)
    if fmt == "docx":
        return count_docx_images(path)
    if fmt == "pptx":
        return count_pptx_images(path)
    raise ValueError("unknown format: %r" % fmt)
